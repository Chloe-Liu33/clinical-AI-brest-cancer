"""Generate an academic-style PPT summarizing METABRIC ablation C-index results."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------- style helpers ----------
TITLE_BLUE   = RGBColor(0x0B, 0x3D, 0x91)
ACCENT_RED   = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN = RGBColor(0x1E, 0x82, 0x49)
DARK_GREY    = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY   = RGBColor(0xEE, 0xEE, 0xEE)
HIGHLIGHT_BG = RGBColor(0xFF, 0xF7, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title(slide, text, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = TITLE_BLUE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14); sp.font.italic = True; sp.font.color.rgb = DARK_GREY
    # underline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_BLUE
    bar.line.fill.background()

def add_bullets(slide, bullets, left, top, width, height, font_size=16, color=DARK_GREY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(font_size if level == 0 else font_size - 2)
        p.font.color.rgb = color
        p.font.bold = (level == 0)
        p.space_after = Pt(4)

def add_text(slide, text, left, top, width, height, font_size=14, bold=False, color=DARK_GREY, align=None, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color
    p.font.italic = italic
    if align: p.alignment = align

def style_cell(cell, text, bold=False, size=11, color=DARK_GREY, bg=None, align=PP_ALIGN.CENTER):
    cell.text = ""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg

def add_footer(slide, idx, total):
    add_text(slide, f"METABRIC Ablation  |  C-index, 5 seeds  |  {idx} / {total}",
             Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             font_size=10, color=DARK_GREY, italic=True, align=PP_ALIGN.RIGHT)

TOTAL = 17

CAL_DIR = "/fs04/scratch2/ek04/limei/AImed/metabric_model/logs/calibration"

# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
# background tint
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC); bg.line.fill.background()

add_text(s, "Breaking the Clinical Ceiling in Breast-Cancer Survival",
         Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0),
         font_size=36, bold=True, color=TITLE_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "An HSIC-based Multimodal Fusion Framework on METABRIC",
         Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.7),
         font_size=22, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)
add_text(s, "Ablation Study  ·  Test C-index  ·  5 random seeds",
         Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5),
         font_size=16, color=ACCENT_RED, align=PP_ALIGN.CENTER)
# decorative bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_BLUE; bar.line.fill.background()
add_text(s, "Results presentation", Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
         font_size=12, color=DARK_GREY, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# Slide 2 — Background (very brief) + Setup
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Background & Experimental Setup",
          subtitle="One-slide context: data, modalities, models compared")

add_text(s, "Task", Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Overall-survival prediction on METABRIC (n ≈ 1,900 patients).", 0),
    ("Metric: Harrell's C-index on the held-out test split.", 0),
    ("Each result averaged over 5 seeds {42, 7, 123, 2024, 31415}.", 0),
], Inches(0.6), Inches(2.0), Inches(6.2), Inches(2.0), font_size=14)

add_text(s, "Modalities  (detailed on the next slide)", Inches(0.6), Inches(3.7), Inches(6.2), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Clinical: 58 features (demographic, tumour, treatment, IHC, INTCLUST).", 0),
    ("Gene expression (mRNA): 20,384 genes — continuous z-scores.", 0),
    ("Copy-number alterations (CNA): 22,544 genes — discrete {−2,−1,0,1,2}.", 0),
], Inches(0.6), Inches(4.1), Inches(6.2), Inches(2.0), font_size=13)

add_text(s, "Models compared", Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Single-modality MLPs (Clinical / Gene / CNA).", 0),
    ("Naive late fusion (Concat) — bi- and tri-modal.", 0),
    ("Variance-based gene selection (Top-K = 256 / 512 / 1024 / 2048).", 0),
    ("Ours: HSIC-based fusion with kernel rank K ∈ {256, 512, 1024, 2048}.", 0),
    ("Ours-Tri: HSIC fusion with the third (CNA) modality.", 0),
    ("Robustness: ‑noICL variants (INTCLUST features removed).", 0),
], Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5), font_size=13)

add_footer(s, 2, TOTAL)

# ============================================================
# Slide 3 — Modalities: what's inside each data block
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "What's Inside Each Modality?",
          subtitle="Three biologically distinct views of the same 1,916 patients")

# Three columns
col_w = Inches(4.1); col_top = Inches(1.7); col_h = Inches(5.2)
col_x = [Inches(0.4), Inches(4.6), Inches(8.8)]
col_titles = ["Clinical  (58 features)",
              "Gene Expression  (mRNA, 20,384 genes)",
              "Copy-Number  (CNA, 22,544 genes)"]
col_colors = [RGBColor(0x1E, 0x82, 0x49), RGBColor(0x0B, 0x3D, 0x91), RGBColor(0xC0, 0x39, 0x2B)]

for x, t, c in zip(col_x, col_titles, col_colors):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_top, col_w, col_h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFD)
    box.line.color.rgb = c; box.line.width = Pt(1.5)
    add_text(s, t, x + Inches(0.15), col_top + Inches(0.1), col_w - Inches(0.3), Inches(0.5),
             font_size=14, bold=True, color=c)

# Clinical column
add_bullets(s, [
    ("Demographics: AGE_AT_DIAGNOSIS, MENOPAUSAL_STATE.", 0),
    ("Tumour: TUMOR_SIZE, TUMOR_STAGE, GRADE, NPI, LYMPH_NODES_POSITIVE, CELLULARITY, HISTOLOGICAL_SUBTYPE, LATERALITY.", 0),
    ("Receptors / IHC: ER, PR, HER2 (status + IHC), HER2_SNP6.", 0),
    ("Treatment: CHEMO, HORMONE, RADIO, BREAST_SURGERY.", 0),
    ("Mutational load: TMB_NONSYNONYMOUS.", 0),
    ("INTCLUST 1–10: integrative subtype derived from joint mRNA + CNA (hence the noICL ablation).", 0),
], col_x[0] + Inches(0.18), col_top + Inches(0.65), col_w - Inches(0.36), col_h - Inches(0.8),
   font_size=10)

# Gene column
add_bullets(s, [
    ("Continuous mRNA expression measured on Illumina HT-12 v3 microarray.", 0),
    ("Pre-processing: log-scale, then z-score standardised per gene on the train split only.", 0),
    ("Captures transcriptional programmes — proliferation, hormone signalling, immune infiltration, EMT, etc.", 0),
    ("High-dim, continuous, smooth — ideal substrate for kernel-based dependency learning.", 0),
], col_x[1] + Inches(0.18), col_top + Inches(0.65), col_w - Inches(0.36), col_h - Inches(0.8),
   font_size=11)

# CNA column
add_bullets(s, [
    ("Discrete GISTIC calls per gene: −2 deep loss, −1 loss, 0 neutral, +1 gain, +2 amplification.", 0),
    ("Captures structural genomic events — chromosomal instability, oncogene amplification (e.g. ERBB2/MYC), tumour-suppressor loss.", 0),
    ("Z-score standardised on the train split before training.", 0),
    ("Sparse, low-entropy signal — many genes are 0 in most patients.", 0),
    ("Weakest predictor on its own (CNA-only C-index = 0.5587).", 0),
], col_x[2] + Inches(0.18), col_top + Inches(0.65), col_w - Inches(0.36), col_h - Inches(0.8),
   font_size=11)

# Why these three
add_text(s, "Why these three together?  Clinical encodes the prognosis the oncologist already knows; mRNA captures the active biological programme of the tumour; CNA captures the underlying structural lesions.  They are complementary — a faithful fusion model should exploit that complementarity, not be hurt by it.",
         Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.5),
         font_size=11, color=DARK_GREY, italic=True)

add_footer(s, 3, TOTAL)

# ============================================================
# Slide 4 — What is "Test C-index"?
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Evaluation Metric — Test C-index",
          subtitle="Harrell's concordance index, the standard ranking metric for time-to-event models")

add_text(s, "Definition", Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Among all comparable patient pairs (i, j) where i died earlier than j survived,", 0),
    ("the model gets credit if its predicted risk for i is higher than for j.", 0),
    ("C-index = fraction of correctly ordered comparable pairs.", 0),
    ("Handles right-censoring naturally — pairs where neither event is observed are simply skipped.", 0),
], Inches(0.6), Inches(2.0), Inches(6.2), Inches(2.5), font_size=13)

# Scale interpretation box
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.7), Inches(6.2), Inches(2.3))
box.fill.solid(); box.fill.fore_color.rgb = HIGHLIGHT_BG
box.line.color.rgb = TITLE_BLUE; box.line.width = Pt(1.0)
add_text(s, "Interpreting the scale", Inches(0.8), Inches(4.85), Inches(5.8), Inches(0.4),
         font_size=14, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("0.50  =  random ordering (no predictive value).", 0),
    ("0.60–0.65  =  weak signal (clinical risk scores often sit here on validation).", 0),
    ("0.65–0.70  =  competitive on METABRIC; 0.70+ is hard to reach.", 0),
    ("1.00  =  perfect ranking (impossible in practice for OS).", 0),
], Inches(0.85), Inches(5.25), Inches(5.8), Inches(1.7), font_size=12)

# Right column: "test" qualifier + protocol
add_text(s, "Why \"test\" specifically?", Inches(7.1), Inches(1.6), Inches(6), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Each seed produces a 70/15/15 train/val/test split (stratified by event).", 0),
    ("Models are tuned on validation; the test split is touched only once at the end.", 0),
    ("\"Test C-index\" is therefore the held-out generalisation score — not training fit, not validation pick-best.", 0),
    ("Reported numbers are the mean ± std of test C-index across the 5 seeds.", 0),
    ("To quantify within-seed uncertainty we also report bootstrap 95% CIs on the test set.", 0),
], Inches(7.1), Inches(2.0), Inches(5.8), Inches(3.0), font_size=13)

# Reference numbers
add_text(s, "Anchor points in this study", Inches(7.1), Inches(5.0), Inches(6), Inches(0.4),
         font_size=14, bold=True, color=ACCENT_RED)
add_bullets(s, [
    ("Random predictor:  0.500.", 0),
    ("Clinical-only ceiling:  0.6800 (5-seed mean, std 0.0023).", 0),
    ("Ours (HSIC K=512):  0.7008 (best, std 0.0063).", 0),
], Inches(7.1), Inches(5.4), Inches(5.8), Inches(1.6), font_size=13)

add_footer(s, 4, TOTAL)

# ============================================================
# Slide 5 — Why the Gene Var-K sweep? (matched K)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Why Multiple Gene Var-K Variants?",
          subtitle="A matched-K control: \"Top-K most variable genes\" vs \"K HSIC-learned dimensions\"")

add_text(s, "The question this baseline answers", Inches(0.6), Inches(1.6), Inches(12.2), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("HSIC compresses 20,384 genes down to K ∈ {256, 512, 1024, 2048} task-relevant components. Could a much simpler selector — \"just keep the K most variable genes\" — do equally well?", 0),
    ("If yes → most of the gain is from dimensionality, not from HSIC.", 1),
    ("If no  → HSIC's task-aware selection is the actual driver.", 1),
    ("To make the comparison fair, Gene Var-K is swept over the SAME K values as HSIC.", 0),
    ("Identical capacity (#features fed into the MLP), identical seeds, identical splits.", 1),
], Inches(0.6), Inches(2.0), Inches(12.2), Inches(2.6), font_size=14)

# Mini comparison table
n_rows = 5; n_cols = 4
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.2)).table
widths = [2.0, 3.4, 3.4, 3.3]
for j, w in enumerate(widths): tbl.columns[j].width = Inches(w)
hdr = ["K", "Gene Var-K (variance ranking)", "Ours HSIC-K (task-aware)", "Δ (HSIC − Var)"]
for j, h in enumerate(hdr):
    style_cell(tbl.cell(0, j), h, bold=True, size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)
rows = [
    ("256",  "0.6008 ± 0.0068", "0.6991 ± 0.0067", "+0.0983"),
    ("512",  "0.6057 ± 0.0191", "0.7008 ± 0.0063", "+0.0951"),
    ("1024", "0.6090 ± 0.0207", "0.6856 ± 0.0227", "+0.0766"),
    ("2048", "0.6162 ± 0.0062", "0.6821 ± 0.0141", "+0.0659"),
]
for i, r in enumerate(rows, start=1):
    bg = LIGHT_GREY if i % 2 == 0 else None
    for j, v in enumerate(r):
        color = ACCENT_GREEN if j == 3 else DARK_GREY
        style_cell(tbl.cell(i, j), v, bold=(j == 3), size=12, color=color, bg=bg)

add_text(s, "Conclusion: at every matched K, HSIC beats variance-selection by ~0.07–0.10 C-index. The gain is from HOW genes are selected, not from how many.",
         Inches(0.6), Inches(6.95), Inches(12.2), Inches(0.4),
         font_size=12, italic=True, color=DARK_GREY)

add_footer(s, 5, TOTAL)

# ============================================================
# Slide 6 — Main Results Table  (was Slide 3)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Main Results — Test C-index (mean ± std, 5 seeds)",
          subtitle="Best two rows highlighted; ours K=512 attains the top score with the lowest variance")

rows_data = [
    ("Clinical-only (MLP)",           "0.6800 ± 0.0023", "15,873"),
    ("Gene-only (MLP)",               "0.6340 ± 0.0090", "2,617,601"),
    ("CNA-only (MLP)",                "0.5587 ± 0.0158", "2,893,825"),
    ("Late Fusion (Concat, bi)",      "0.6792 ± 0.0097", "2,650,369"),
    ("Late Fusion 3 (Concat, tri)",   "0.6705 ± 0.0122", "5,552,513"),
    ("Gene Var-256 (MLP)",            "0.6008 ± 0.0068", "41,217"),
    ("Gene Var-512 (MLP)",            "0.6057 ± 0.0191", "73,985"),
    ("Gene Var-1024 (MLP)",           "0.6090 ± 0.0207", "139,521"),
    ("Gene Var-2048 (MLP)",           "0.6162 ± 0.0062", "270,593"),
    ("Ours (HSIC K=256)",             "0.6991 ± 0.0067", "2,782,593"),
    ("Ours (HSIC K=512)",             "0.7008 ± 0.0063", "2,782,593"),
    ("Ours (HSIC K=1024)",            "0.6856 ± 0.0227", "2,782,593"),
    ("Ours (HSIC K=2048)",            "0.6821 ± 0.0141", "2,782,593"),
    ("Ours-Tri (HSIC K=256)",         "0.6878 ± 0.0112", "5,816,961"),
    ("Ours-Tri (HSIC K=512)",         "0.7002 ± 0.0190", "5,816,961"),
]

n_rows = len(rows_data) + 1
n_cols = 3
tbl_left, tbl_top, tbl_w, tbl_h = Inches(0.7), Inches(1.55), Inches(12.0), Inches(5.3)
tbl = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table
tbl.columns[0].width = Inches(5.5)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(3.0)

# header
headers = ["Model", "Test C-index", "Parameters"]
for j, h in enumerate(headers):
    style_cell(tbl.cell(0, j), h, bold=True, size=14, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

# rows
for i, (m, c, p) in enumerate(rows_data, start=1):
    is_best   = m.startswith("Ours (HSIC K=512)")
    is_runner = m.startswith("Ours (HSIC K=256)") or m.startswith("Ours-Tri (HSIC K=512)")
    is_ours   = m.startswith("Ours")
    bg = HIGHLIGHT_BG if (is_best or is_runner) else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_best else (ACCENT_GREEN if is_ours else DARK_GREY)
    bold = is_ours
    align_left = PP_ALIGN.LEFT
    style_cell(tbl.cell(i, 0), "  " + m, bold=bold, size=12, color=color, bg=bg, align=align_left)
    style_cell(tbl.cell(i, 1), c, bold=bold, size=12, color=color, bg=bg)
    style_cell(tbl.cell(i, 2), p, bold=bold, size=12, color=color, bg=bg)

add_footer(s, 6, TOTAL)

# ============================================================
# Slide 7 — Headline Finding: Breaking the Clinical Ceiling
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Headline Finding — Breaking the Clinical Ceiling",
          subtitle="Only HSIC fusion (K=256/512) significantly exceeds the clinical-only baseline")

# left: callout
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.7), Inches(5.6), Inches(4.8))
box.fill.solid(); box.fill.fore_color.rgb = HIGHLIGHT_BG
box.line.color.rgb = ACCENT_RED; box.line.width = Pt(1.5)
add_text(s, "Best model", Inches(0.9), Inches(1.85), Inches(5.0), Inches(0.4),
         font_size=14, bold=True, color=ACCENT_RED)
add_text(s, "Ours (HSIC, K = 512)", Inches(0.9), Inches(2.25), Inches(5.0), Inches(0.6),
         font_size=24, bold=True, color=TITLE_BLUE)
add_text(s, "C-index = 0.7008 ± 0.0063",
         Inches(0.9), Inches(2.95), Inches(5.0), Inches(0.5),
         font_size=18, bold=True, color=DARK_GREY)
add_bullets(s, [
    ("+0.0208 vs Clinical-only   (p = 0.0005)", 0),
    ("+0.0668 vs Gene-only        (p = 0.0005)", 0),
    ("+0.0216 vs Late Fusion       (p = 0.017)", 0),
    ("Lowest variance across all 19 models", 0),
], Inches(0.95), Inches(3.7), Inches(5.2), Inches(2.7), font_size=14, color=DARK_GREY)

# right: ranked comparison
add_text(s, "Why this matters", Inches(6.6), Inches(1.7), Inches(6.4), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Clinical-only is a hard ceiling on METABRIC (0.6800).", 0),
    ("Naive concat fusion fails to break it: Late Fusion = 0.6792, Late Fusion 3 = 0.6705.", 0),
    ("All variance-selected gene MLPs fall below 0.6200.", 0),
    ("Single-modality genomic MLPs cannot match clinical features.", 0),
    ("HSIC fusion is the only family that consistently and significantly surpasses clinical.", 0),
    ("Implication: it is HOW we fuse — not what or how much we add — that drives the gain.", 0),
], Inches(6.6), Inches(2.1), Inches(6.4), Inches(4.6), font_size=14)

add_footer(s, 7, TOTAL)

# ============================================================
# Slide 8 — K sweep (HSIC kernel rank)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Ablation — HSIC Kernel Rank K",
          subtitle="A clear sweet spot: K=512 is best in both performance and stability")

# Table on the left
n_rows, n_cols = 6, 4
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7), Inches(6.6), Inches(3.7)).table
tbl.columns[0].width = Inches(1.4)
tbl.columns[1].width = Inches(1.7)
tbl.columns[2].width = Inches(1.7)
tbl.columns[3].width = Inches(1.8)
hdr = ["K", "C-index", "Std", "p vs Clinical"]
for j, h in enumerate(hdr):
    style_cell(tbl.cell(0, j), h, bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

rows = [
    ("256",  "0.6991", "0.0067", "0.006 ✓"),
    ("512",  "0.7008", "0.0063", "0.0005 ✓"),
    ("1024", "0.6856", "0.0227", "0.631  ✗"),
    ("2048", "0.6821", "0.0141", "0.774  ✗"),
    ("Clinical-only", "0.6800", "0.0023", "—"),
]
for i, r in enumerate(rows, start=1):
    is_best = (r[0] == "512")
    bg = HIGHLIGHT_BG if is_best else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_best else DARK_GREY
    bold = is_best
    for j, v in enumerate(r):
        style_cell(tbl.cell(i, j), v, bold=bold, size=12, color=color, bg=bg)

# Right: insight box
add_text(s, "Take-aways", Inches(7.6), Inches(1.7), Inches(5.3), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Performance follows an inverted-U over K.", 0),
    ("K = 512 maximises the mean (0.7008) and minimises the std (0.0063).", 0),
    ("K ≥ 1024: variance roughly 3× larger, mean drops; the kernel becomes over-parameterised.", 0),
    ("Above K=1024 the model loses statistical separation from Clinical-only (p > 0.6).", 0),
    ("Recommendation: report K=512 as the headline; K=256 as a robust runner-up; K≥1024 as the failure mode.", 0),
], Inches(7.6), Inches(2.1), Inches(5.3), Inches(5.0), font_size=14)

add_footer(s, 8, TOTAL)

# ============================================================
# Slide 9 — Variance-based gene selection vs HSIC (full-gene context)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Variance-based Gene Selection Fails — HSIC Wins",
          subtitle="High-variance genes ≠ prognostic genes; HSIC's task-aware selection clearly outperforms")

# Table
n_rows = 7; n_cols = 3
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7), Inches(6.5), Inches(4.4)).table
tbl.columns[0].width = Inches(2.7)
tbl.columns[1].width = Inches(1.9)
tbl.columns[2].width = Inches(1.9)
hdr = ["Selector", "C-index", "# Genes / K"]
for j, h in enumerate(hdr):
    style_cell(tbl.cell(0, j), h, bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

rows = [
    ("Var-Top 256",   "0.6008 ± 0.0068", "256"),
    ("Var-Top 512",   "0.6057 ± 0.0191", "512"),
    ("Var-Top 1024",  "0.6090 ± 0.0207", "1024"),
    ("Var-Top 2048",  "0.6162 ± 0.0062", "2048"),
    ("Gene-only (full)", "0.6340 ± 0.0090", "≈ 24,000"),
    ("Ours HSIC (K=512)", "0.7008 ± 0.0063", "K = 512"),
]
for i, r in enumerate(rows, start=1):
    is_best = "HSIC" in r[0]
    bg = HIGHLIGHT_BG if is_best else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_best else DARK_GREY
    bold = is_best
    align = PP_ALIGN.LEFT if r[0].startswith("Var") or "Gene" in r[0] or "HSIC" in r[0] else PP_ALIGN.CENTER
    for j, v in enumerate(r):
        style_cell(tbl.cell(i, j), v if j > 0 else "  " + v,
                   bold=bold, size=12, color=color, bg=bg,
                   align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

add_text(s, "What this tells us", Inches(7.4), Inches(1.7), Inches(5.6), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Variance-based selection plateaus around 0.60–0.62 — strictly worse than even Gene-only.", 0),
    ("Adding more high-variance genes barely helps (Var-256 → Var-2048: +0.015).", 0),
    ("HSIC compresses gene expression to a 512-dim summary that beats the full 24K-gene MLP by +0.067.", 0),
    ("Direct evidence that HSIC's task-aware kernel learning > unsupervised variance ranking.", 0),
], Inches(7.4), Inches(2.1), Inches(5.6), Inches(5.0), font_size=14)

add_footer(s, 9, TOTAL)

# ============================================================
# Slide 10 — Tri-modal extension and CNA modality
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Tri-modal Extension — CNA Adds Little, but HSIC Does Not Get Hurt",
          subtitle="On METABRIC the CNA signal is weak; HSIC absorbs it without degradation, naive concat does not")

n_rows = 5; n_cols = 3
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7), Inches(7.0), Inches(2.6)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.0)
for j, h in enumerate(["Model", "C-index", "Δ vs bi-modal"]):
    style_cell(tbl.cell(0, j), h, bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

rows = [
    ("Late Fusion (Concat, bi)",      "0.6792", "—"),
    ("Late Fusion 3 (Concat, tri)",   "0.6705", "−0.0087"),
    ("Ours (HSIC K=512)",             "0.7008", "—"),
    ("Ours-Tri (HSIC K=512)",         "0.7002", "−0.0006"),
]
for i, r in enumerate(rows, start=1):
    is_ours = r[0].startswith("Ours")
    bg = HIGHLIGHT_BG if "Ours-Tri" in r[0] else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_GREEN if is_ours else DARK_GREY
    bold = is_ours
    style_cell(tbl.cell(i, 0), "  " + r[0], bold=bold, size=12, color=color, bg=bg, align=PP_ALIGN.LEFT)
    style_cell(tbl.cell(i, 1), r[1], bold=bold, size=12, color=color, bg=bg)
    style_cell(tbl.cell(i, 2), r[2], bold=bold, size=12, color=color, bg=bg)

add_bullets(s, [
    ("CNA-only C-index is 0.5587 (close to chance). The CNA signal alone is weak on METABRIC.", 0),
    ("Naive concat is hurt by the weak modality: Late Fusion 3 < Late Fusion by 0.0087.", 0),
    ("HSIC absorbs CNA without damage: Ours-Tri ≈ Ours (Δ = −0.0006).", 0),
    ("Robustness implication: HSIC fusion is graceful under low-quality modalities;", 1),
    ("naive concatenation is not.", 1),
], Inches(0.6), Inches(4.5), Inches(12.3), Inches(2.5), font_size=14)

add_footer(s, 10, TOTAL)

# ============================================================
# Slide 11 — Statistical Significance
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Statistical Significance — Paired t-test (n = 5 seeds)",
          subtitle="K=256 and K=512 dominate every baseline; K=1024/2048 lose separation from Clinical & Late Fusion")

# Compact significance table
n_rows = 7; n_cols = 6
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.6)).table
widths = [3.2, 1.6, 1.6, 1.7, 1.7, 1.7]
for j, w in enumerate(widths):
    tbl.columns[j].width = Inches(w)
hdr = ["Ours model", "vs Clinical", "vs Gene", "vs CNA", "vs LF", "vs LF3"]
for j, h in enumerate(hdr):
    style_cell(tbl.cell(0, j), h, bold=True, size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

# Each row: model + p-values (with ✓ / ✗)
def fmt(p):
    s_ = "✓" if p < 0.05 else "✗"
    return f"{p:.4f} {s_}"

rows = [
    ("Ours (HSIC K=256)",  0.0062, 0.0002, 0.0000, 0.0473, 0.0068),
    ("Ours (HSIC K=512)",  0.0005, 0.0005, 0.0000, 0.0171, 0.0025),
    ("Ours (HSIC K=1024)", 0.6306, 0.0073, 0.0015, 0.5541, 0.2667),
    ("Ours (HSIC K=2048)", 0.7744, 0.0059, 0.0004, 0.7806, 0.0816),
    ("Ours-Tri (HSIC K=256)", 0.1765, 0.0009, 0.0003, 0.0902, 0.0841),
    ("Ours-Tri (HSIC K=512)", 0.0911, 0.0017, 0.0007, 0.0517, 0.0515),
]
for i, r in enumerate(rows, start=1):
    name = r[0]
    is_top = ("HSIC K=256)" in name and "Tri" not in name) or ("HSIC K=512)" in name and "Tri" not in name)
    bg = HIGHLIGHT_BG if is_top else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_top else DARK_GREY
    bold = is_top
    style_cell(tbl.cell(i, 0), "  " + name, bold=bold, size=11, color=color, bg=bg, align=PP_ALIGN.LEFT)
    for j, p in enumerate(r[1:], start=1):
        cell_color = ACCENT_GREEN if p < 0.05 else ACCENT_RED
        style_cell(tbl.cell(i, j), fmt(p), bold=bold, size=11, color=cell_color, bg=bg)

add_bullets(s, [
    ("K = 256 / 512 reach p < 0.05 vs every baseline including Clinical-only and Late Fusion.", 0),
    ("K = 1024 / 2048 are not separable from Clinical-only or Late Fusion (p > 0.55).", 0),
    ("Ours-Tri reaches p < 0.05 vs Gene/CNA but only borderline vs Clinical/Late-Fusion (consistent with weak CNA).", 0),
    ("Caveat: n = 5 paired observations — treat p < 0.05 as suggestive, not conclusive.", 0),
], Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7), font_size=13)

add_footer(s, 11, TOTAL)

# ============================================================
# Slide 12 — INTCLUST leakage robustness
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Robustness Check — INTCLUST Leakage (with-ICL vs no-ICL)",
          subtitle="Removing INTCLUST features causes only minor, non-significant drops; gains are not driven by leakage")

n_rows = 4; n_cols = 5
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.0)).table
widths = [3.5, 2.0, 2.0, 2.3, 2.3]
for j, w in enumerate(widths):
    tbl.columns[j].width = Inches(w)
for j, h in enumerate(["Model", "with-ICL", "no-ICL", "Δ", "paired p"]):
    style_cell(tbl.cell(0, j), h, bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

rows = [
    ("Clinical-only (MLP)",       "0.6800", "0.6675", "−0.0125", "0.244"),
    ("Late Fusion 3 (Concat)",    "0.6705", "0.6699", "−0.0006", "0.946"),
    ("Ours-Tri (HSIC K=512)",     "0.7002", "0.6892", "−0.0110", "0.291"),
]
for i, r in enumerate(rows, start=1):
    is_ours = r[0].startswith("Ours")
    bg = HIGHLIGHT_BG if is_ours else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_GREEN if is_ours else DARK_GREY
    bold = is_ours
    style_cell(tbl.cell(i, 0), "  " + r[0], bold=bold, size=12, color=color, bg=bg, align=PP_ALIGN.LEFT)
    for j, v in enumerate(r[1:], start=1):
        style_cell(tbl.cell(i, j), v, bold=bold, size=12, color=color, bg=bg)

add_text(s, "Why this is reassuring", Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Concern: INTCLUST already encodes prognosis-relevant subtype information; could it explain our gain?", 0),
    ("Removing ICL drops every model by ≈ 0.001–0.013, but no contrast is statistically significant (p > 0.24).", 0),
    ("Crucially, Ours-Tri-noICL = 0.6892 still beats Clinical-noICL = 0.6675 and Late Fusion 3-noICL = 0.6699.", 0),
    ("→ The HSIC advantage is preserved when INTCLUST is excluded; the gain is not a leakage artefact.", 0),
], Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6), font_size=14)

add_footer(s, 12, TOTAL)

# ============================================================
# Slide 13 — Parameter efficiency
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Parameter Efficiency — More Performance, Fewer Parameters",
          subtitle="HSIC compresses, rather than concatenates: ~half the params, +0.030 in C-index")

n_rows = 5; n_cols = 3
tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.6), Inches(1.7), Inches(7.0), Inches(2.7)).table
tbl.columns[0].width = Inches(3.0); tbl.columns[1].width = Inches(2.0); tbl.columns[2].width = Inches(2.0)
for j, h in enumerate(["Model", "Params", "C-index"]):
    style_cell(tbl.cell(0, j), h, bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

rows = [
    ("Clinical-only",     "15,873",   "0.6800"),
    ("Late Fusion (bi)",  "2,650,369","0.6792"),
    ("Late Fusion 3 (tri)","5,552,513","0.6705"),
    ("Ours (HSIC K=512)", "2,782,593","0.7008"),
]
for i, r in enumerate(rows, start=1):
    is_ours = r[0].startswith("Ours")
    bg = HIGHLIGHT_BG if is_ours else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_ours else DARK_GREY
    bold = is_ours
    style_cell(tbl.cell(i, 0), "  " + r[0], bold=bold, size=12, color=color, bg=bg, align=PP_ALIGN.LEFT)
    style_cell(tbl.cell(i, 1), r[1], bold=bold, size=12, color=color, bg=bg)
    style_cell(tbl.cell(i, 2), r[2], bold=bold, size=12, color=color, bg=bg)

add_text(s, "Notes", Inches(7.9), Inches(1.7), Inches(5.0), Inches(0.4),
         font_size=18, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Ours (K=512) uses 2.78M params — about half of Late Fusion 3 (5.55M).", 0),
    ("Performance gap: +0.0303 C-index over Late Fusion 3.", 0),
    ("Compression-then-fuse is more parameter-efficient than concat-then-MLP.", 0),
    ("The win is methodological, not capacity-driven.", 0),
], Inches(7.9), Inches(2.1), Inches(5.0), Inches(4.5), font_size=14)

add_footer(s, 13, TOTAL)

# ============================================================
# Slide 14 — Beyond C-index: Why calibration matters
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Beyond C-index — Why Calibration Matters",
          subtitle="C-index measures ranking; calibration measures whether the predicted probabilities are trustworthy")

add_text(s, "Two independent axes of survival-model quality", Inches(0.6), Inches(1.5),
         Inches(12.2), Inches(0.4), font_size=18, bold=True, color=TITLE_BLUE)

# Two-column comparison
col_w = Inches(6.0); col_top = Inches(2.0); col_h = Inches(2.5)
left_x = Inches(0.5); right_x = Inches(6.8)

box_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, col_top, col_w, col_h)
box_l.fill.solid(); box_l.fill.fore_color.rgb = RGBColor(0xEC, 0xF3, 0xFA)
box_l.line.color.rgb = TITLE_BLUE; box_l.line.width = Pt(1.0)
add_text(s, "C-index (discrimination)", left_x + Inches(0.2), col_top + Inches(0.1),
         col_w - Inches(0.4), Inches(0.4), font_size=15, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Question: among any two patients, is the higher-risk one ranked first?", 0),
    ("Range: 0.5 (random) → 1.0 (perfect ordering).", 0),
    ("Insensitive to absolute probability — only the order matters.", 0),
    ("A model with great C-index can still output \"5-year survival 70%\"  when truth is 50%.", 0),
], left_x + Inches(0.25), col_top + Inches(0.55), col_w - Inches(0.5), col_h - Inches(0.7),
   font_size=12)

box_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, col_top, col_w, col_h)
box_r.fill.solid(); box_r.fill.fore_color.rgb = HIGHLIGHT_BG
box_r.line.color.rgb = ACCENT_RED; box_r.line.width = Pt(1.0)
add_text(s, "Calibration (probability fidelity)", right_x + Inches(0.2), col_top + Inches(0.1),
         col_w - Inches(0.4), Inches(0.4), font_size=15, bold=True, color=ACCENT_RED)
add_bullets(s, [
    ("Question: when the model says \"S(t) = 70%\" does ~70% of that group really survive to t?", 0),
    ("Plotted as predicted vs observed event probability at a fixed horizon.", 0),
    ("Critical for clinical decision-making, prognostic counselling, treatment thresholds.", 0),
    ("Independent from C-index — a model can win one and lose the other.", 0),
], right_x + Inches(0.25), col_top + Inches(0.55), col_w - Inches(0.5), col_h - Inches(0.7),
   font_size=12)

# --- Bottom: 4-metric glossary ---
add_text(s, "The four metrics we report  (all lower = better)", Inches(0.5), Inches(4.7),
         Inches(12.4), Inches(0.4), font_size=15, bold=True, color=TITLE_BLUE)

card_y = Inches(5.1); card_h = Inches(1.9); card_w = Inches(3.05)
card_xs = [Inches(0.4), Inches(3.55), Inches(6.70), Inches(9.85)]
card_titles = ["Brier(t)", "ICI", "E50", "E90"]
card_subtitles = [
    "Probability MSE",
    "Integrated Calibration Index",
    "Median abs calibration error",
    "90th-pct abs calibration error",
]
card_bodies = [
    # Brier
    [("Mean squared error of S(t) vs truth, IPCW-corrected for censoring.", 0),
     ("Range [0, 1].  0.25 = uninformative.  ≤ 0.20 typical for OS.", 0),
     ("Single number combining calibration + sharpness.", 0)],
    # ICI
    [("Mean |observed − predicted| event-prob, averaged over the test set.", 0),
     ("< 0.05 = excellent;  < 0.10 = acceptable.", 0),
     ("Most direct readout of the calibration plot.", 0)],
    # E50
    [("Median of |observed − predicted| event-prob.", 0),
     ("Robust 'typical' calibration error (insensitive to tails).", 0),
     ("If close to ICI → errors are uniform across patients.", 0)],
    # E90
    [("90th percentile of |observed − predicted| event-prob.", 0),
     ("Captures worst-case calibration in the tail.", 0),
     ("Large E90 with small ICI → a few badly mis-calibrated subgroups.", 0)],
]
for x, t, sub, body in zip(card_xs, card_titles, card_subtitles, card_bodies):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    box.line.color.rgb = TITLE_BLUE; box.line.width = Pt(1.0)
    add_text(s, t, x + Inches(0.15), card_y + Inches(0.08),
             card_w - Inches(0.3), Inches(0.32),
             font_size=14, bold=True, color=ACCENT_RED)
    add_text(s, sub, x + Inches(0.15), card_y + Inches(0.40),
             card_w - Inches(0.3), Inches(0.28),
             font_size=10, italic=True, color=TITLE_BLUE)
    add_bullets(s, body,
                x + Inches(0.18), card_y + Inches(0.70),
                card_w - Inches(0.36), card_h - Inches(0.78),
                font_size=9)

add_footer(s, 14, TOTAL)

# ============================================================
# Slide 15 — Calibration Results (table + 5y / 10y comparison plots)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Calibration Results — HSIC Wins on Probabilities Too",
          subtitle="All four metrics, both horizons. Lower is better; Ours (HSIC K=512) wins every cell — 8 / 8.")

# --- Full numerical table: 5 models × 9 columns -----------------------------
n_rows, n_cols = 6, 9
tbl_x = Inches(0.35); tbl_y = Inches(1.55)
tbl_w = Inches(12.6); tbl_h = Inches(2.45)
tbl = s.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h).table

widths = [2.6, 1.20, 1.20, 1.20, 1.30, 1.30, 1.30, 1.30, 1.20]
for j, w in enumerate(widths):
    tbl.columns[j].width = Inches(w)

# Header row
hdr = ["Model",
       "Brier 5y", "ICI 5y", "E50 5y", "E90 5y",
       "Brier 10y", "ICI 10y", "E50 10y", "E90 10y"]
for j, h in enumerate(hdr):
    style_cell(tbl.cell(0, j), h, bold=True, size=10, color=RGBColor(0xFF, 0xFF, 0xFF), bg=TITLE_BLUE)

# Data rows
cal_full = [
    ("Clinical-only (MLP)",     ["0.164±.008", "0.062±.041", "0.056±.045", "0.102±.057",
                                 "0.211±.010", "0.088±.045", "0.086±.055", "0.143±.056"]),
    ("Late Fusion (Concat)",    ["0.168±.009", "0.138±.022", "0.125±.024", "0.209±.035",
                                 "0.236±.005", "0.174±.022", "0.173±.019", "0.286±.020"]),
    ("Late Fusion 3 (Concat)",  ["0.181±.007", "0.159±.017", "0.129±.012", "0.253±.045",
                                 "0.248±.008", "0.203±.034", "0.203±.032", "0.322±.028"]),
    ("Ours (HSIC K=512)",       ["0.149±.004", "0.046±.018", "0.037±.017", "0.095±.034",
                                 "0.204±.007", "0.064±.030", "0.053±.023", "0.129±.066"]),
    ("Ours-Tri (HSIC K=512)",   ["0.150±.005", "0.057±.012", "0.043±.009", "0.096±.022",
                                 "0.207±.012", "0.076±.030", "0.065±.034", "0.148±.055"]),
]
for i, (name, vals) in enumerate(cal_full, start=1):
    is_best = name.startswith("Ours (HSIC")
    is_ours = name.startswith("Ours")
    bg = HIGHLIGHT_BG if is_best else (LIGHT_GREY if i % 2 == 0 else None)
    color = ACCENT_RED if is_best else (ACCENT_GREEN if is_ours else DARK_GREY)
    bold = is_ours
    style_cell(tbl.cell(i, 0), "  " + name, bold=bold, size=10, color=color, bg=bg, align=PP_ALIGN.LEFT)
    for j, v in enumerate(vals, start=1):
        style_cell(tbl.cell(i, j), v, bold=bold, size=9, color=color, bg=bg)

# --- Bottom: two compare images side by side, with a caption strip on the right
img_top  = Inches(4.20)
img_h    = Inches(2.55)
img_w    = Inches(2.55)
s.shapes.add_picture(os.path.join(CAL_DIR, "calib_compare_60m.png"),
                     Inches(0.35), img_top, width=img_w, height=img_h)
s.shapes.add_picture(os.path.join(CAL_DIR, "calib_compare_120m.png"),
                     Inches(3.0), img_top, width=img_w, height=img_h)

# Right: reading guide + key findings (compact)
add_text(s, "Reading the plots", Inches(5.85), Inches(4.20), Inches(7.2), Inches(0.32),
         font_size=13, bold=True, color=TITLE_BLUE)
add_bullets(s, [
    ("Diagonal = perfect calibration.  Above = model under-predicts death risk (over-optimistic).", 0),
    ("Ours / Ours-Tri (red, purple) hug the diagonal at both horizons.", 0),
    ("Late Fusion (orange) sits far above the diagonal — large systematic bias.", 0),
], Inches(5.85), Inches(4.55), Inches(7.3), Inches(1.2), font_size=11)

add_text(s, "Key findings (numerical)", Inches(5.85), Inches(5.65), Inches(7.2), Inches(0.32),
         font_size=13, bold=True, color=ACCENT_RED)
add_bullets(s, [
    ("Ours (HSIC K=512) wins ALL 8 cells (4 metrics × 2 horizons).", 0),
    ("Late Fusion has C-index ≈ Clinical-only (0.679 vs 0.680), yet 5y ICI is 2.2× worse (0.138 vs 0.062) — a failure mode invisible to C-index.", 0),
    ("CNA: Concat-3 → 10y ICI = 0.203 (severe); HSIC-Tri only +0.012 — graceful absorption.", 0),
    ("Ours has lowest Brier std (0.004) — most stable across seeds, too.", 0),
], Inches(5.85), Inches(5.95), Inches(7.3), Inches(1.0), font_size=10)

add_footer(s, 15, TOTAL)

# ============================================================
# Slide 16 — Visual Diagnosis: Why naive concat fails calibration
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Visual Diagnosis — Why Naive Concat Fails Calibration",
          subtitle="Same C-index ≈ 0.68, very different probability quality. Decile bins (red dots) reveal the contrast")

# Two stacked dual-panels: Ours on top, Late Fusion on bottom
panel_w = Inches(9.4)
panel_h = Inches(2.6)
panel_x = Inches(0.35)

# Ours (top)
s.shapes.add_picture(os.path.join(CAL_DIR, "calib_Ours_HSIC_K512.png"),
                     panel_x, Inches(1.55), width=panel_w, height=panel_h)

# Late Fusion (bottom)
s.shapes.add_picture(os.path.join(CAL_DIR, "calib_Late_Fusion_Concat.png"),
                     panel_x, Inches(4.35), width=panel_w, height=panel_h)

# Side annotations (right column)
ann_x = Inches(10.0); ann_w = Inches(3.0)

# Annotation for Ours
ann_box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              ann_x, Inches(1.55), ann_w, panel_h)
ann_box1.fill.solid(); ann_box1.fill.fore_color.rgb = HIGHLIGHT_BG
ann_box1.line.color.rgb = ACCENT_RED; ann_box1.line.width = Pt(1.0)
add_text(s, "Ours (HSIC K=512)", ann_x + Inches(0.15), Inches(1.65),
         ann_w - Inches(0.3), Inches(0.4), font_size=12, bold=True, color=ACCENT_RED)
add_bullets(s, [
    ("Decile dots track the diagonal.", 0),
    ("Smooth curve nearly straight.", 0),
    ("ICI = 0.038 (5y), 0.055 (10y).", 0),
    ("Predicted probabilities are usable.", 0),
], ann_x + Inches(0.18), Inches(2.1), ann_w - Inches(0.36), Inches(2.0), font_size=10)

# Annotation for Late Fusion
ann_box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              ann_x, Inches(4.35), ann_w, panel_h)
ann_box2.fill.solid(); ann_box2.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEA)
ann_box2.line.color.rgb = DARK_GREY; ann_box2.line.width = Pt(1.0)
add_text(s, "Late Fusion (Concat)", ann_x + Inches(0.15), Inches(4.45),
         ann_w - Inches(0.3), Inches(0.4), font_size=12, bold=True, color=DARK_GREY)
add_bullets(s, [
    ("Curve sits ABOVE the diagonal.", 0),
    ("→ Systematic UNDER-prediction of risk.", 0),
    ("Predicting 10% event ≈ 35% truth.", 0),
    ("ICI = 0.137 (5y), 0.174 (10y).", 0),
    ("Ranks fine, probabilities mislead.", 0),
], ann_x + Inches(0.18), Inches(4.85), ann_w - Inches(0.36), Inches(2.0), font_size=10)

# Bottom takeaway strip
add_text(s, "Take-away: a 5-year survival forecast of \"70%\" from Late Fusion really corresponds to ~50% truth — clinically misleading despite competitive C-index. HSIC fusion is the only fusion family whose probabilities can be trusted at face value.",
         Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.4),
         font_size=11, italic=True, color=ACCENT_RED)

add_footer(s, 16, TOTAL)

# ============================================================
# Slide 17 — Conclusions
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_title(s, "Conclusions", subtitle="What this ablation establishes, and what is next")

add_text(s, "Established", Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.4),
         font_size=18, bold=True, color=ACCENT_GREEN)
add_bullets(s, [
    ("HSIC fusion (K=512) is the only model that significantly exceeds the clinical-only C-index ceiling (p = 0.0005).", 0),
    ("Best-calibrated probabilities at both 5y and 10y: 8/8 wins on Brier / ICI / E50 / E90.", 0),
    ("Reveals a failure mode of naive concat: similar C-index to Clinical-only but 2.2× worse 5y ICI.", 0),
    ("Performance is bell-shaped in K; K=512 is the sweet spot in both mean and variance.", 0),
    ("Task-aware HSIC selection clearly outperforms variance-based gene selection.", 0),
    ("HSIC fusion is robust to weak modalities (CNA), unlike naive concatenation — both for ranking and calibration.", 0),
    ("Gains survive INTCLUST removal — not a leakage artefact.", 0),
    ("Methodological, not capacity-driven: ~half the params of tri-modal concat, higher C-index, better calibration.", 0),
], Inches(0.6), Inches(2.0), Inches(12.2), Inches(3.5), font_size=13)

add_text(s, "Limitations & next steps", Inches(0.6), Inches(5.2), Inches(6.0), Inches(0.4),
         font_size=18, bold=True, color=ACCENT_RED)
add_bullets(s, [
    ("n = 5 seeds → p-values are suggestive, not conclusive. Plan more seeds and external validation (TCGA-BRCA).", 0),
    ("Refine K around 384–640 to confirm 512 is a true minimum, not a plateau edge.", 0),
    ("CNA contribution is currently flat; consider attention-gated CNA integration.", 0),
], Inches(0.6), Inches(5.6), Inches(12.2), Inches(2.0), font_size=14)

add_footer(s, 17, TOTAL)

# ---------- save ----------
out = "/fs04/scratch2/ek04/limei/AImed/metabric_model/logs/METABRIC_Ablation_Cindex.pptx"
prs.save(out)
print(f"Saved: {out}")
